#!/usr/bin/env python3
"""
Generate an editable .pptx Marketing Wins slide from weekly data.
Output imports into Google Slides with full text editing capability.

Usage:
    python scripts/slide_to_pptx.py          # uses hardcoded data below
    # Future: accept JSON input for dynamic weeks
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand colors ──
BLUE = RGBColor(0x2C, 0x64, 0xE3)
DARK_BLUE = RGBColor(0x0F, 0x1C, 0x47)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_BG = RGBColor(0xF9, 0xFA, 0xFB)
BORDER_GRAY = RGBColor(0xE5, 0xE7, 0xEB)
CALLOUT_BG = RGBColor(0xEF, 0xF6, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Funnel colors
GREEN = RGBColor(0x4A, 0xDE, 0x80)
TEAL = RGBColor(0x2D, 0xD4, 0xBF)
RED = RGBColor(0xF8, 0x71, 0x71)
ORANGE = RGBColor(0xFB, 0x92, 0x3C)
LIGHT_BLUE = RGBColor(0x60, 0xA5, 0xFA)

# ── Slide dimensions (16:9) ──
SLIDE_W = Inches(13.333)  # 1920px at 144dpi
SLIDE_H = Inches(7.5)     # 1080px at 144dpi

# ── Helpers ──

def add_textbox(slide, left, top, width, height, text, font_size=12,
                bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                font_name="Arial"):
    """Add a simple text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rich_textbox(slide, left, top, width, height, runs, alignment=PP_ALIGN.LEFT):
    """Add a text box with multiple formatted runs.
    runs = [(text, font_size, bold, color), ...]
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    for i, (text, font_size, bold, color) in enumerate(runs):
        if i == 0:
            run = p.runs[0] if p.runs else p.add_run()
            run.text = text
        else:
            run = p.add_run()
            run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Arial"
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color=LIGHT_BG,
                     border_color=BORDER_GRAY, no_border=False):
    """Add a rounded rectangle shape with subtle border."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if no_border:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = border_color
        shape.line.width = Emu(6350)  # 0.5pt — very thin
    shape.adjustments[0] = 0.04
    return shape


def add_metric_card(slide, left, top, width, height, value, label,
                    value_color=BLACK):
    """Add a metric card with value + label."""
    add_rounded_rect(slide, left, top, width, height)
    # Value
    add_textbox(slide, left, top + Inches(0.08), width, Inches(0.45),
                str(value), font_size=28, bold=True, color=value_color,
                alignment=PP_ALIGN.CENTER)
    # Label
    add_textbox(slide, left, top + Inches(0.52), width, Inches(0.3),
                label, font_size=8, bold=False, color=GRAY,
                alignment=PP_ALIGN.CENTER)


def add_section_header(slide, left, top, width, text):
    """Add a blue section header with underline."""
    tb = add_textbox(slide, left, top, width, Inches(0.35), text,
                     font_size=16, bold=True, color=BLUE)
    # Blue underline
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top + Inches(0.32), width, Pt(2)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()
    return tb


def add_funnel_bar(slide, left, top, bar_width, height, text, color):
    """Add a colored funnel bar with text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, bar_width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.adjustments[0] = 0.15
    # Text inside bar
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Arial"


def add_email_row(slide, left, top, width, name, stats_text):
    """Add an email campaign row."""
    row_h = Inches(0.5)
    add_rounded_rect(slide, left, top, width, row_h)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.04), width - Inches(0.2), Inches(0.22),
                name, font_size=10, bold=True, color=BLACK)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.25), width - Inches(0.2), Inches(0.22),
                stats_text, font_size=9, bold=False, color=GRAY)


def add_bullet_list(slide, left, top, width, items, font_size=11):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, Inches(len(items) * 0.28))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = BLACK
        p.font.name = "Arial"
        p.space_after = Pt(4)
        p.level = 0
        # Blue bullet
        pPr = p._pPr
        if pPr is None:
            from lxml import etree
            nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            pPr = etree.SubElement(p._p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
        from lxml import etree
        nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        buFont = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buFont')
        buFont.set('typeface', 'Arial')
        buChar = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
        buChar.set('char', '\u2022')
        buClr = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buClr')
        srgb = etree.SubElement(buClr, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        srgb.set('val', '2C64E3')
    return txBox


# ══════════════════════════════════════════════════════
#  WEEKLY DATA — Edit this section each week
# ══════════════════════════════════════════════════════

WEEK_LABEL = "Week of Feb 28 – Mar 6, 2026"

INBOUND_METRICS = [
    ("8", "Contact Sales\n(deduped)"),
    ("4", "Sales Outreach\nCompleted"),
    ("2", "Meetings\nScheduled"),
    ("93", "Login\n(deduped)"),
    ("74", "Verification Help\n(deduped)"),
]

FUNNEL = [
    ("Form Fill", 8, 1.0, GREEN),
    ("Assigned to Rep", 8, 1.0, TEAL),
    ("First Contact", 4, 0.5, RED),
    ("Meeting (SAL)", 2, 0.25, ORANGE),
    ("Deal Created", 1, 0.125, LIGHT_BLUE),
]

EMAIL_CAMPAIGNS = [
    ("FCM Webinar Invites (Main List)", "Delivered: ~2,046  |  Opens: ~22.9%  |  CTO: 58.1%"),
    ("FCM Webinar Closed-Lost Invite", "Delivered: 667  |  Opens: 35.7%  |  CTO: 54.2%"),
    ("CMG Case Study Email", "Delivered: 3,761  |  Opens: 23.1%  |  CTO: 58.3%"),
]

AD_SPEND = [
    # (platform, spend, impressions, cpm, clicks, ctr, cpc, leads)
    ("Meta", "$93", "5,255", "$17.74", "44", "0.84%", "$2.12", "0"),
    ("LinkedIn", "$450", "3,520", "$127.84", "18", "0.51%", "$25.00", "1"),
]
AD_SPEND_TOTAL = ("Total", "$543", "8,775", "$61.91", "62", "0.71%", "—", "1")
AD_SPEND_PERIOD = "Mar 2–5"

WEBINAR_METRICS = [
    ("22", "Registrants", BLUE),
    ("2,720", "Emails Delivered", BLACK),
]

WEBINAR_CALLOUT = (
    "Stood up entire campaign this week — project plan, audience segments, "
    "Knock/Wrapper, Clay table, 2 invite sends, LinkedIn post, plus closed-lost "
    "re-engagement to 494 contacts"
)

SEGMENTATION = [
    ("Found 39,956 untagged contacts",
     " (19.1% of CRM) — built auto-tagging script that classified 33,333 business emails by vertical. Tagged 4,273 companies automatically."),
    ("New lending audience: 3,301 contacts",
     " — queried 5 lending verticals (IMB, Bank, CU, Home Equity, Subprime) with 12-month engagement signals. Synced to Knock for campaign sends."),
]

HUB_ONELINER = "Bi-directional HubSpot ↔ Linear sync. Read-only dashboard for the whole team — no access needed. truv-brain.vercel.app/hub"

UPCOMING = [
    "FCM Webinar Invites 3 & 4 + day-of reminder",
    "Product Update Email — live send",
    "BankSouth Case Study Email — live send",
    "FCM Closed-Lost reminder emails (2 of 3, 3 of 3)",
    "Case Study emails to Mortgage Ops (CCM, PHM, FCM)",
    "Whitepaper Email Campaign send",
    "March Insider Public Sector Webinar setup",
]


# ══════════════════════════════════════════════════════
#  BUILD THE SLIDE
# ══════════════════════════════════════════════════════

def build_slide():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Blank layout
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)

    # ── Column positions ──
    L_LEFT = Inches(0.5)
    L_WIDTH = Inches(6.8)
    R_LEFT = Inches(7.6)
    R_WIDTH = Inches(5.4)

    # ══ HEADER ══
    add_textbox(slide, L_LEFT, Inches(0.25), Inches(5), Inches(0.6),
                "Marketing Wins", font_size=36, bold=True, color=BLUE)
    add_textbox(slide, L_LEFT, Inches(0.78), Inches(4), Inches(0.3),
                WEEK_LABEL, font_size=11, color=GRAY)

    # ══ LEFT COLUMN ══
    y = Inches(1.2)

    # Inbound header
    add_section_header(slide, L_LEFT, y, L_WIDTH, "Inbound")
    y += Inches(0.5)

    # Metric cards
    card_w = Inches(1.26)
    card_h = Inches(0.85)
    card_gap = Inches(0.1)
    for i, (val, label) in enumerate(INBOUND_METRICS):
        cx = L_LEFT + i * (card_w + card_gap)
        add_metric_card(slide, cx, y, card_w, card_h, val, label)
    y += Inches(1.05)

    # Conversion Funnel
    add_section_header(slide, L_LEFT, y, L_WIDTH, "Conversion Funnel")
    y += Inches(0.5)

    funnel_label_w = Inches(1.2)
    funnel_max_w = Inches(5.5)
    bar_h = Inches(0.28)
    for label, count, pct, color in FUNNEL:
        add_textbox(slide, L_LEFT, y, funnel_label_w, bar_h,
                    label, font_size=9, color=GRAY, alignment=PP_ALIGN.RIGHT)
        bar_w = max(Inches(0.6), Emu(int(funnel_max_w * pct)))
        add_funnel_bar(slide, L_LEFT + funnel_label_w + Inches(0.12), y,
                       bar_w, bar_h, count, color)
        y += Inches(0.33)
    y += Inches(0.1)

    # Email Campaigns Sent
    add_section_header(slide, L_LEFT, y, L_WIDTH, "Email Campaigns Sent")
    y += Inches(0.5)

    for name, stats in EMAIL_CAMPAIGNS:
        add_email_row(slide, L_LEFT, y, L_WIDTH, name, stats)
        y += Inches(0.55)
    y += Inches(0.05)

    # Paid Media section
    add_section_header(slide, L_LEFT, y, L_WIDTH, f"Paid Media ({AD_SPEND_PERIOD})")
    y += Inches(0.45)

    # Table header
    col_widths = [Inches(0.95), Inches(0.7), Inches(0.85), Inches(0.85), Inches(0.7), Inches(0.7), Inches(0.75), Inches(0.7)]
    headers = ["Platform", "Spend", "Impr.", "CPM", "Clicks", "CTR", "CPC", "Leads"]
    x_pos = L_LEFT
    for hw, ht in zip(col_widths, headers):
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos, y, hw, Inches(0.28))
        rect.fill.solid()
        rect.fill.fore_color.rgb = LIGHT_BG
        rect.line.fill.background()
        add_textbox(slide, x_pos, y + Inches(0.02), hw, Inches(0.24),
                    ht, font_size=7, bold=True, color=DARK_BLUE,
                    alignment=PP_ALIGN.CENTER)
        x_pos += hw
    y += Inches(0.3)

    # Table rows
    all_rows = AD_SPEND + [AD_SPEND_TOTAL]
    for row_data in all_rows:
        is_total = (row_data == AD_SPEND_TOTAL)
        x_pos = L_LEFT
        for i, (cw, val) in enumerate(zip(col_widths, row_data)):
            if is_total:
                rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos, y, cw, Inches(0.26))
                rect.fill.solid()
                rect.fill.fore_color.rgb = LIGHT_BG
                rect.line.fill.background()
            add_textbox(slide, x_pos, y + Inches(0.01), cw, Inches(0.24),
                        val, font_size=9, bold=is_total, color=BLACK,
                        alignment=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER)
            x_pos += cw
        y += Inches(0.28)

    # ══ RIGHT COLUMN ══
    ry = Inches(1.2)

    # FCM Webinar header
    add_section_header(slide, R_LEFT, ry, R_WIDTH, "FCM Webinar")
    ry += Inches(0.5)

    # Webinar metric cards
    wm_w = Inches(2.6)
    wm_h = Inches(0.85)
    for i, (val, label, vcolor) in enumerate(WEBINAR_METRICS):
        wx = R_LEFT + i * (wm_w + Inches(0.1))
        add_metric_card(slide, wx, ry, wm_w, wm_h, val, label,
                        value_color=vcolor)
    ry += Inches(1.0)

    # Webinar callout
    callout_h = Inches(0.65)
    add_rounded_rect(slide, R_LEFT, ry, R_WIDTH, callout_h,
                     fill_color=CALLOUT_BG, no_border=True)
    # Blue left border
    border = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, R_LEFT, ry, Pt(3), callout_h
    )
    border.fill.solid()
    border.fill.fore_color.rgb = BLUE
    border.line.fill.background()

    add_rich_textbox(slide, R_LEFT + Inches(0.15), ry + Inches(0.06),
                     R_WIDTH - Inches(0.3), callout_h - Inches(0.1),
                     [
                         ("Stood up entire campaign this week", 9, True, BLUE),
                         (" — " + WEBINAR_CALLOUT.split("— ", 1)[1], 9, False, BLACK),
                     ])
    ry += Inches(0.8)

    # Segmentation & Targeting
    add_section_header(slide, R_LEFT, ry, R_WIDTH, "Segmentation & Targeting")
    ry += Inches(0.5)

    for bold_text, rest_text in SEGMENTATION:
        callout_h = Inches(0.6)
        add_rounded_rect(slide, R_LEFT, ry, R_WIDTH, callout_h,
                         fill_color=CALLOUT_BG, no_border=True)
        border = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, R_LEFT, ry, Pt(3), callout_h
        )
        border.fill.solid()
        border.fill.fore_color.rgb = BLUE
        border.line.fill.background()
        add_rich_textbox(slide, R_LEFT + Inches(0.15), ry + Inches(0.06),
                         R_WIDTH - Inches(0.3), callout_h - Inches(0.1),
                         [
                             (bold_text, 9, True, BLUE),
                             (rest_text, 9, False, BLACK),
                         ])
        ry += Inches(0.68)

    # Marketing Hub one-liner
    add_section_header(slide, R_LEFT, ry, R_WIDTH, "What's New: Truv Marketing Hub")
    ry += Inches(0.45)
    add_textbox(slide, R_LEFT, ry, R_WIDTH, Inches(0.35),
                HUB_ONELINER, font_size=9, color=GRAY)
    ry += Inches(0.45)

    # Upcoming Next Week
    add_section_header(slide, R_LEFT, ry, R_WIDTH, "Upcoming Next Week")
    ry += Inches(0.5)

    add_bullet_list(slide, R_LEFT, ry, R_WIDTH, UPCOMING, font_size=11)

    # ── Save ──
    out_dir = os.path.join(os.path.dirname(__file__), "..", "docs", "presentations")
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, "marketing-wins-2026-03-06.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}")
    return out_path


if __name__ == "__main__":
    build_slide()
